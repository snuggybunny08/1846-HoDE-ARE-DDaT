"""
Generate PowerPoint presentation for Agriculture Data Platform Architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import io
from PIL import Image

def create_presentation():
    """Create the PowerPoint presentation for the platform design."""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Agricultural Data Platform"
    subtitle.text = "Modern Cloud-Native Architecture for ARE Directorate\n\nHead of Data Engineering\nTechnical Assessment\nNovember 2025"
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    
    content2 = slide2.placeholders[1]
    tf = content2.text_frame
    
    p1 = tf.add_paragraph()
    p1.text = "Vision"
    p1.font.bold = True
    p1.level = 0
    
    p2 = tf.add_paragraph()
    p2.text = "Transform ARE's data capabilities with a modern, scalable Unified Data Platform (UDP)"
    p2.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "\nKey Objectives"
    p3.font.bold = True
    p3.level = 0
    
    p4 = tf.add_paragraph()
    p4.text = "• Replace legacy systems with cloud-native architecture"
    p4.level = 1
    
    p5 = tf.add_paragraph()
    p5.text = "• Enable real-time analytics for £600M annual payments"
    p5.level = 1
    
    p6 = tf.add_paragraph()
    p6.text = "• Support IoT integration from agricultural sensors"
    p6.level = 1
    
    p7 = tf.add_paragraph()
    p7.text = "• Ensure GDPR compliance and data governance"
    p7.level = 1
    
    # Slide 3: Design Principles
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Core Design Principles"
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    
    principles = [
        ("Scalability", "Auto-scaling infrastructure to handle seasonal peaks (harvest periods)"),
        ("Security", "Zero-trust architecture with encryption at rest and in transit"),
        ("Interoperability", "API-first design supporting REST, GraphQL, and event streaming"),
        ("Governance", "Built-in data lineage, quality monitoring, and GDPR compliance"),
        ("Cost Efficiency", "Serverless-first approach with consumption-based pricing"),
        ("Resilience", "Multi-AZ deployment with 99.99% availability SLA")
    ]
    
    for principle, description in principles:
        p = tf3.add_paragraph()
        p.text = f"• {principle}: {description}"
        p.level = 0
        p.font.size = Pt(14)
    
    # Slide 4: Platform Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide4.shapes.title.text = "Platform Architecture - Layered Approach"
    
    # Add architecture diagram description
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    textbox = slide4.shapes.add_textbox(left, top, width, height)
    tf4 = textbox.text_frame
    
    layers = [
        ("Data Sources", "• Farm IoT Sensors • Legacy Systems • Spreadsheets • External APIs", "Blue"),
        ("Ingestion Layer", "• Azure Event Hubs (streaming) • Azure Data Factory (batch) • API Gateway", "Light Blue"),
        ("Storage Layer", "• Bronze (Raw): Azure Data Lake Gen2 • Silver (Cleansed): Delta Lake • Gold (Curated): Synapse", "Green"),
        ("Processing Layer", "• Apache Spark (Databricks) • Stream Analytics • Azure Functions", "Orange"),
        ("Analytics Layer", "• Power BI • Azure ML • Cognitive Services", "Purple"),
        ("Governance Layer", "• Purview (Catalog) • Data Quality Rules • Lineage Tracking", "Red")
    ]
    
    for layer_name, components, color in layers:
        p = tf4.add_paragraph()
        p.text = f"{layer_name}"
        p.font.bold = True
        p.font.size = Pt(14)
        
        p2 = tf4.add_paragraph()
        p2.text = components
        p2.font.size = Pt(12)
        p2.level = 1
        
    # Slide 5: Technology Stack
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Technology Stack & Patterns"
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    
    tech_stack = [
        ("Cloud Platform", "Microsoft Azure (aligning with Scottish Government standards)"),
        ("Data Lake", "Azure Data Lake Storage Gen2 with Delta Lake format"),
        ("Orchestration", "Azure Data Factory for batch, Event Hubs for streaming"),
        ("Processing", "Databricks for large-scale processing, Functions for microservices"),
        ("Data Warehouse", "Azure Synapse Analytics for enterprise reporting"),
        ("Metadata", "Azure Purview for data catalog and governance"),
        ("Monitoring", "Azure Monitor, Application Insights, Log Analytics"),
        ("Security", "Azure Key Vault, Managed Identities, Private Endpoints")
    ]
    
    for tech, description in tech_stack:
        p = tf5.add_paragraph()
        p.text = f"• {tech}: {description}"
        p.font.size = Pt(12)
    
    # Slide 6: Data Quality & Lineage
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Data Quality & Lineage Framework"
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    
    quality_items = [
        "Automated Quality Checks",
        "• Completeness, Uniqueness, Consistency, Validity, Timeliness",
        "",
        "Data Lineage Tracking",
        "• Source-to-target mapping with Azure Purview",
        "• Impact analysis for schema changes",
        "• Audit trail for compliance (GDPR Article 30)",
        "",
        "Metadata Management",
        "• Business glossary and data dictionary",
        "• Schema versioning with Delta Lake",
        "• Automated documentation generation"
    ]
    
    for item in quality_items:
        p = tf6.add_paragraph()
        p.text = item
        if not item.startswith("•") and item:
            p.font.bold = True
        p.font.size = Pt(12)
    
    # Slide 7: Delivery Approach
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Delivery Approach & Timeline"
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    
    phases = [
        "Phase 1: Foundation (Months 1-3)",
        "• Setup cloud infrastructure and DevOps pipelines",
        "• Establish data governance framework",
        "• Proof of Concept with sensor data",
        "",
        "Phase 2: Migration (Months 4-9)",
        "• Migrate priority data sources to Data Lake",
        "• Build core data pipelines and quality checks",
        "• Deploy analytics dashboards for pilot users",
        "",
        "Phase 3: Scale (Months 10-12)",
        "• Full production deployment",
        "• Machine learning model deployment",
        "• Training and handover to operations team"
    ]
    
    for phase in phases:
        p = tf7.add_paragraph()
        p.text = phase
        if "Phase" in phase:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 32, 96)
        p.font.size = Pt(12)
    
    # Save presentation
    prs.save('AgricultureDataPlatform_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
